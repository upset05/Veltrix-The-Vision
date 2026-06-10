import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Light Theme)
    BG_COLOR = RGBColor(248, 250, 252) # slate-50
    PRIMARY = RGBColor(79, 70, 229)    # indigo-600
    TEXT_DARK = RGBColor(15, 23, 42)    # slate-900
    TEXT_MUTED = RGBColor(71, 85, 105)  # slate-600
    WHITE = RGBColor(255, 255, 255)
    
    PURPLE = RGBColor(139, 92, 246)
    EMERALD = RGBColor(16, 185, 129)
    ORANGE = RGBColor(249, 115, 22)

    # Helper function to set slide background
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper to add standard slide header
    def add_header(slide, title_text, category_text):
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.name = 'Arial'
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK

    blank_slide_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Minimalist & Premium)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide1)
    
    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "VELTRIX ECOSYSTEM"
    p.font.name = 'Arial'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "The Vision & Architecture"
    p2.font.name = 'Arial'
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_DARK
    
    # Subtitle / Creator info
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.5))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    p_info = tf_info.paragraphs[0]
    p_info.text = "A Unified Suite of Fintech Platforms and Intelligent Assistants\nCreated by The Accountant · Version by Version"
    p_info.font.name = 'Arial'
    p_info.font.size = Pt(14)
    p_info.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: Ecosystem Overview
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide2)
    add_header(slide2, "The Connected Digital Ecosystem", "Overview")
    
    # 4 Columns for the 4 Stages
    col_width = Inches(2.7)
    gap = Inches(0.3)
    start_left = Inches(0.8)
    top_pos = Inches(2.0)
    height = Inches(4.5)
    
    stages = [
        {"num": "01", "name": "Veltrix VTU", "status": "BUILT / OPERATIONAL", "color": EMERALD, "desc": "Premium student fintech platform. Allows instant purchases of data & airtime, secure wallet funding via Paystack, transaction tracking, and admin controls."},
        {"num": "02", "name": "Veltrix Core", "status": "ACTIVE DEVELOPMENT", "color": PURPLE, "desc": "AI-powered intelligence system. Contains modular files for learning, identity databases, local facts memory, session reports, notes and tasks tracker."},
        {"num": "03", "name": "Veltrix Studio", "status": "PLANNED", "color": PRIMARY, "desc": "AI website builder. Generates landing pages, copy, templates, layouts, and custom deployment structures. Will integrate closely with Veltrix Core."},
        {"num": "04", "name": "Veltrix OS", "status": "LONG-TERM VISION", "color": ORANGE, "desc": "Android-based custom mobile OS. Houses Veltrix Core as the underlying intelligence layer for AI-assisted phone experience and supervised tasks."}
    ]
    
    for i, st in enumerate(stages):
        left_pos = start_left + i * (col_width + gap)
        
        # Add box shape for background card
        shape = slide2.shapes.add_shape(
            1, # Rectangle
            left_pos, top_pos, col_width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(226, 232, 240) # border slate-200
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.3)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        
        # Stage number
        p_num = tf.paragraphs[0]
        p_num.text = f"STAGE {st['num']}"
        p_num.font.name = 'Arial'
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = st['color']
        
        # Title
        p_title = tf.add_paragraph()
        p_title.text = st['name']
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.space_before = Pt(8)
        p_title.space_after = Pt(4)
        
        # Status Label
        p_status = tf.add_paragraph()
        p_status.text = st['status']
        p_status.font.name = 'Arial'
        p_status.font.size = Pt(9)
        p_status.font.bold = True
        p_status.font.color.rgb = st['color']
        p_status.space_after = Pt(14)
        
        # Description
        p_desc = tf.add_paragraph()
        p_desc.text = st['desc']
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 3: Veltrix VTU (Operational)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide3)
    add_header(slide3, "Veltrix VTU: Campus Airtime & Data", "Stage 01 — Operational")
    
    # Left side: Details
    info_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Built for Campus Life"
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(14)
    
    bullets = [
        "Instant Top-Ups: Airtime and data delivered in milliseconds via premium provider APIs.",
        "Secure Payment Flow: Seamless wallet funding secured by Paystack flat-fee transactions.",
        "Security-First: User login, persistent session storage, and 4-digit transaction PIN verification.",
        "Referral Engine: Cash bonuses instantly credited to user wallets on referral funding.",
        "Management Hub: Full dashboard reporting, transaction logs, and administrative controls."
    ]
    for b in bullets:
        p_bullet = tf.add_paragraph()
        p_bullet.text = "• " + b
        p_bullet.font.name = 'Arial'
        p_bullet.font.size = Pt(12)
        p_bullet.font.color.rgb = TEXT_MUTED
        p_bullet.space_after = Pt(8)

    # Right side: Add image (Homepage screenshot)
    img_path = os.path.join(os.path.dirname(__file__), 'assets', 'vtu-dashboard.png')
    if os.path.exists(img_path):
        slide3.shapes.add_picture(img_path, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))

    # -------------------------------------------------------------
    # SLIDE 4: Veltrix Core (Evolution Timeline)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide4)
    add_header(slide4, "Veltrix Core: Visual Release Evolution", "Stage 02 — Active Development")
    
    # Add timeline nodes
    timeline_y = Inches(2.5)
    node_w = Inches(1.8)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    timeline_steps = [
        ("v0.1", "Genesis", "Pulse Dot"),
        ("v0.2", "Awakening", "Active Panel"),
        ("v0.4", "Memory", "Storage Recall"),
        ("v0.8", "Modular", "Code Refactor"),
        ("v0.95", "Productive", "Notes & Tasks"),
        ("v1.0", "Dream Alpha", "Official Release")
    ]
    
    for i, step in enumerate(timeline_steps):
        left_pos = start_left + i * (node_w + gap)
        
        # Add card box
        shape = slide4.shapes.add_shape(
            1, left_pos, timeline_y, node_w, Inches(3.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        
        p_v = tf.paragraphs[0]
        p_v.text = step[0]
        p_v.font.name = 'Arial'
        p_v.font.size = Pt(22)
        p_v.font.bold = True
        p_v.font.color.rgb = PURPLE
        p_v.alignment = PP_ALIGN.CENTER
        
        p_title = tf.add_paragraph()
        p_title.text = step[1]
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(12)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_before = Pt(6)
        
        p_desc = tf.add_paragraph()
        p_desc.text = step[2]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 5: Veltrix Core Architecture
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide5)
    add_header(slide5, "Veltrix Core: Internal System Modules", "Stage 02 — Architecture")
    
    # 3 Column layout
    col_width = Inches(3.6)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(2.2)
    height = Inches(4.2)
    
    modules = [
        {"title": "Identity & Knowledge", "desc": "Reads configuration dynamically from dedicated JSON database categories (anime, personal, business focus). Stores state and local identities and updates categories via prompt instructions."},
        {"title": "Memory & Fact Engine", "desc": "Implements persistent session tracking using localStorage. Learns facts dynamically via prompt scripts, index lookups, exports local data reports, and removes forgotten records."},
        {"title": "Productivity Utilities", "desc": "Integrates user-focused note collections, interactive lists of tasks with toggleable completion checkboxes, aliases, and keyword search indexes across everything stored."}
    ]
    
    for i, mod in enumerate(modules):
        left_pos = start_left + i * (col_width + gap)
        
        shape = slide5.shapes.add_shape(
            1, left_pos, top_pos, col_width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(226, 232, 240)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.3)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        
        p_title = tf.paragraphs[0]
        p_title.text = mod["title"]
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = PURPLE
        p_title.space_after = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = mod["desc"]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 6: Veltrix Core Long-Term Vision
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide6)
    add_header(slide6, "Veltrix Core: The Long-Term Vision", "Stage 02 — Core Roadmap")
    
    # Left side: Details
    info_box = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "A Supervised Ambient Assistant"
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(14)
    
    visions = [
        "Dynamic Learning: Learn continuously from connected documents, logs, databases, and APIs.",
        "Active Research: Access the internet to query details, summarize reports, and verify sources.",
        "Workspace Interaction: Read, analyze, write, and organize files in standard project directories.",
        "Self-Improvement: Suggest logic tweaks and improvements to its own engine dynamically.",
        "Supervisor Guardrails: Access computer-level tools only under direct human supervision.",
        "Strict User Approval: Never run major actions, modify files, or execute builds without explicit user permission."
    ]
    for v in visions:
        p_bullet = tf.add_paragraph()
        p_bullet.text = "• " + v
        p_bullet.font.name = 'Arial'
        p_bullet.font.size = Pt(12)
        p_bullet.font.color.rgb = TEXT_MUTED
        p_bullet.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 7: Future Stages: Veltrix Studio & Veltrix OS
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide7)
    add_header(slide7, "Future Frameworks: Studio & OS", "Stage 03 & 04 — Expansion")
    
    # Left column: Veltrix Studio
    left_shape = slide7.shapes.add_shape(
        1, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5)
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = WHITE
    left_shape.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_left = left_shape.text_frame
    tf_left.word_wrap = True
    tf_left.margin_top = Inches(0.3)
    tf_left.margin_left = Inches(0.2)
    tf_left.margin_right = Inches(0.2)
    
    p = tf_left.paragraphs[0]
    p.text = "Veltrix Studio"
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(10)
    
    studio_points = [
        "Status: Planned / Future",
        "AI Website Generation: Instantly builds clean landing pages, portfolios, and blogs.",
        "Smart Content Copy: Automated layout design and text copy tailored for businesses.",
        "Interactive Customization: Drag-and-drop structural edits with deployable source files.",
        "Integrated Engine: Future direct control mapping via connected Veltrix Core assistants."
    ]
    for sp in studio_points:
        p_point = tf_left.add_paragraph()
        p_point.text = "• " + sp
        p_point.font.name = 'Arial'
        p_point.font.size = Pt(11)
        p_point.font.color.rgb = TEXT_MUTED
        p_point.space_after = Pt(6)

    # Right column: Veltrix OS
    right_shape = slide7.shapes.add_shape(
        1, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5)
    )
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = WHITE
    right_shape.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_right = right_shape.text_frame
    tf_right.word_wrap = True
    tf_right.margin_top = Inches(0.3)
    tf_right.margin_left = Inches(0.2)
    tf_right.margin_right = Inches(0.2)
    
    p_r = tf_right.paragraphs[0]
    p_r.text = "Veltrix OS"
    p_r.font.name = 'Arial'
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = ORANGE
    p_r.space_after = Pt(10)
    
    os_points = [
        "Status: Long-Term Vision",
        "Android-Based Platform: A custom OS built strictly for Android mobile devices.",
        "Ambient Layer: Veltrix Core functions as the central system intelligence.",
        "App Automation: Supervised AI assistant automates screen actions, checks inputs, and coordinates workflows.",
        "Approval Protection: Built with user guardrails, preventing system tasks from running without explicit notification."
    ]
    for op in os_points:
        p_point = tf_right.add_paragraph()
        p_point.text = "• " + op
        p_point.font.name = 'Arial'
        p_point.font.size = Pt(11)
        p_point.font.color.rgb = TEXT_MUTED
        p_point.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 8: Closing Quote (Even Dots Can Dream)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide8)
    
    # Big quote in center
    quote_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.0))
    tf = quote_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "“Even dots can dream.”"
    p.font.name = 'Arial'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.italic = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(24)
    
    p2 = tf.add_paragraph()
    p2.text = "From VTU infrastructure to intelligent systems, Veltrix is being built version by version into a connected digital ecosystem."
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), 'Veltrix_the_Vision_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    create_presentation()
